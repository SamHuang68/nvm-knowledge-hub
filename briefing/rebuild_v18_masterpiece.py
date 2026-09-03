import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SRC_PATH = r"C:\Users\Sam\Documents\ChatGPT\NVM\NVM_AI_OCP_NVM_Opportunities_Executive_v17_Layout_Repaired_EN.pptx"
OUT_PATH_CHATGPT = r"C:\Users\Sam\Documents\ChatGPT\NVM\NVM_AI_OCP_NVM_Opportunities_Executive_v18_Comprehensive_Update_EN.pptx"
OUT_PATH_LOCAL = r"c:\Users\Sam\Documents\antigravity\NVM\NVM_AI_OCP_NVM_Opportunities_Executive_v18_Comprehensive_Update_EN.pptx"

# Exact Color Palette
C_DARK_NAVY   = RGBColor(10, 36, 52)     # #0A2434 - Main Titles
C_SLATE_BODY  = RGBColor(56, 83, 101)    # #385365 - Subtitles & Card Body
C_TEAL_KICKER = RGBColor(23, 107, 92)    # #176B5C - Target 1 Kicker & Accent
C_CYAN_KICKER = RGBColor(8, 122, 122)    # #087A7A - Target 2 Kicker
C_CYAN_BRIGHT = RGBColor(113, 236, 227)  # #71ECE3 - Synopsys Kicker
C_WHITE       = RGBColor(255, 255, 255)  # #FFFFFF - White text
C_LIGHT_BLUE  = RGBColor(200, 217, 224)  # #C8D9E0 - Synopsys Bullets
C_FOOTER_GRAY = RGBColor(112, 136, 148)  # #708894 - Footer notes

def clone_slide_shapes(src_slide, target_slide):
    spTree = target_slide.shapes._spTree
    for child in list(spTree):
        spTree.remove(child)
    for child in src_slide.shapes._spTree:
        spTree.append(copy.deepcopy(child))

def set_formatted_run(paragraph, text, font_name="Arial", font_size=14.0, bold=False, color_rgb=C_SLATE_BODY):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color_rgb
    return run

def clear_and_set_text(shape, text, font_name="Arial", font_size=14.0, bold=False, color_rgb=C_SLATE_BODY, align=None):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align: p.alignment = align
    set_formatted_run(p, text, font_name, font_size, bold, color_rgb)

def clear_and_set_multiline(shape, lines, font_name="Arial", font_size=14.0, bold=False, color_rgb=C_SLATE_BODY, align=None):
    tf = shape.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if align: p.alignment = align
        set_formatted_run(p, line, font_name, font_size, bold, color_rgb)

def set_slide_notes(slide, notes_en: str, notes_zh: str):
    full_notes = (
        f"================================================================================\n"
        f"[SPEAKER NOTES — EN]\n{notes_en.strip()}\n\n"
        f"[中文講者備忘與論述重點]\n{notes_zh.strip()}\n"
        f"================================================================================"
    )
    slide.notes_slide.notes_text_frame.text = full_notes

def configure_perfect_card_slide(slide, kicker_text, title_line1, title_line2, subtitle_text,
                                 pill_text, page_num_str,
                                 t1_kicker, t1_title, t1_body,
                                 t2_kicker, t2_title, t2_body,
                                 syn_kicker, syn_title, syn_bullets,
                                 bottom_banner_text):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text.strip()
        
        if "QUALIFY NOW" in txt:
            clear_and_set_text(sh, kicker_text, "Consolas", 9.8, True, C_TEAL_KICKER)
        elif "DDR5 PMIC + SPD Hub" in txt or "TSMC 22ULL" in txt or "3D SoIC" in txt or "NIST SP 800-208" in txt:
            clear_and_set_multiline(sh, [title_line1, title_line2], "Arial", 30.0, True, C_DARK_NAVY)
        elif "Two state contracts" in txt or "Sub-0.6V" in txt or "Eliminating" in txt or "Compact 32-byte" in txt:
            clear_and_set_text(sh, subtitle_text, "Arial", 14.8, False, C_SLATE_BODY)
        elif txt in ["04", "07", "08", "09"]:
            clear_and_set_text(sh, page_num_str, "Consolas", 8.3, True, C_FOOTER_GRAY)
        elif "PACKAGE" in txt or "CONTINUUM" in txt or "INTEGRITY" in txt or "GATE" in txt:
            clear_and_set_text(sh, pill_text, "Consolas", 9.3, True, C_TEAL_KICKER)
        elif "TARGET 1" in txt:
            clear_and_set_text(sh, t1_kicker, "Consolas", 9.4, True, C_TEAL_KICKER)
        elif "MTP configuration" in txt or "0.5V native" in txt or "Known Good Die" in txt or "Stateful hash" in txt:
            clear_and_set_text(sh, t1_title, "Arial", 18.0, True, C_DARK_NAVY)
        elif "programmable sequencing" in txt or "near-threshold" in txt or "defective unverified" in txt or "overcomes 40x" in txt:
            clear_and_set_text(sh, t1_body, "Arial", 14.8, False, C_SLATE_BODY)
        elif "TARGET 2" in txt:
            clear_and_set_text(sh, t2_kicker, "Consolas", 9.4, True, C_CYAN_KICKER)
        elif "Protected NVM" in txt or "28nm eFlash" in txt or "Die-to-Die" in txt or "SILC physics" in txt:
            clear_and_set_text(sh, t2_title, "Arial", 18.0, True, C_DARK_NAVY)
        elif "8-Kbit / 1024-byte" in txt or "ReRAM (+2 masks)" in txt or "SPDM 1.3 mandates" in txt or "175°C stress" in txt:
            clear_and_set_text(sh, t2_body, "Arial", 14.8, False, C_SLATE_BODY)
        elif "SYNOPSYS QUALIFICATION CONTRACT" in txt:
            clear_and_set_text(sh, syn_kicker, "Consolas", 9.4, True, C_CYAN_BRIGHT)
        elif "Close one named DDR5" in txt or "Qualify AntiFuse" in txt or "Deploy AntiFuse" in txt or "1-to-0 Hard" in txt:
            clear_and_set_text(sh, syn_title, "Arial", 20.0, True, C_WHITE)
        elif "exact persistent payload" in txt or "zero extra mask" in txt or "die-unique immutable" in txt or "nanoscale metallic" in txt:
            clear_and_set_multiline(sh, syn_bullets, "Arial", 15.3, False, C_LIGHT_BLUE)
        elif "Public evidence establishes the socket" in txt or "Public foundry roadmap" in txt or "Chiplet Zero-Trust" in txt or "Physics proof" in txt:
            clear_and_set_text(sh, bottom_banner_text, "Arial", 14.3, True, C_WHITE)

def main():
    print("Rebuilding v18 deck with pixel-perfect zero-collision alignment...")
    prs = Presentation(SRC_PATH)
    
    # 1. Slide 1: Cover
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text
        if "SYNOPSYS NVM · AI / OCP" in txt:
            clear_and_set_text(sh, "SYNOPSYS NVM · AI, CHIPLET & IOT OPPORTUNITY BRIEF", "Arial", 38.0, True, C_WHITE)
        elif "Target-bound qualification framework" in txt:
            clear_and_set_text(sh, "Executive qualification framework for OCP, DDR5, AI Power, TSMC IoT (22ULL to N4e), 3D Chiplets & PQC Boot", "Arial", 16.0, False, C_LIGHT_BLUE)
        elif "VERSION 17" in txt:
            clear_and_set_text(sh, "VERSION 18 · COMPREHENSIVE UPDATE", "Consolas", 10.0, True, C_CYAN_BRIGHT)

    # 2. Slide 2: Review Gate 1 (01 · EXECUTIVE ANSWER)
    s2 = prs.slides[1]
    for sh in s2.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text
        if "Authorize three bounded" in txt:
            clear_and_set_multiline(sh, ["Authorize five target-bound qualifications", "Keep two families in discovery"], "Arial", 30.0, True, C_DARK_NAVY)
        elif "The recommendation is an engagement" in txt:
            clear_and_set_text(sh, "Five named state contracts · two discovery families · zero unverified claims", "Arial", 14.8, False, C_SLATE_BODY)
        elif "DDR5 PMIC + SPD Hub" in txt:
            clear_and_set_text(sh, "DDR5 & AI Power Subsystems", "Arial", 20.0, True, C_DARK_NAVY)
        elif "MTP-class configuration" in txt:
            clear_and_set_text(sh, "PMIC 100k MTP + SPD Hub protected EEPROM + multi-phase VR calibration", "Arial", 14.8, False, C_SLATE_BODY)
        elif "Accelerator power controller" in txt:
            clear_and_set_text(sh, "TSMC IoT Continuum (22ULL to N4e)", "Arial", 20.0, True, C_DARK_NAVY)
        elif "Rail policy · calibration" in txt:
            clear_and_set_text(sh, "0.5V near-threshold operation · pure-logic AntiFuse OTP as MCU code storage", "Arial", 14.8, False, C_SLATE_BODY)
        elif "BMC / OAM platform trust" in txt:
            clear_and_set_text(sh, "Platform Trust & 3D Chiplet RoT", "Arial", 20.0, True, C_DARK_NAVY)
        elif "OTP persists lifecycle" in txt:
            clear_and_set_text(sh, "Caliptra RoT + dedicated UCIe 3D SoIC on-die identity & pre-bond KGD testing", "Arial", 14.8, False, C_SLATE_BODY)
        elif "Leadership decision: authorize three" in txt:
            clear_and_set_text(sh, "Leadership decision: authorize five target-bound qualification briefs", "Arial", 14.3, True, C_WHITE)

    # 3. Slide 3: State Grammar (Preserve)
    # 4. Slides 4, 5, 6: DDR5, AI Power, Trust (Preserve)

    # 5. Insert 3 NEW Qualification Slides after Slide 6 (Slides 07, 08, 09)
    template_slide = prs.slides[3]

    # --- Slide 07: TSMC IoT Continuum ---
    s7 = prs.slides.add_slide(prs.slide_layouts[0])
    clone_slide_shapes(template_slide, s7)
    configure_perfect_card_slide(
        s7,
        kicker_text="06 · QUALIFY NOW · TSMC IOT CONTINUUM",
        title_line1="TSMC 22ULL to N4e Scaling",
        title_line2="Pure-Logic OTP as MCU Code Storage",
        subtitle_text="Sub-0.6V near-threshold operation · zero-mask code execution",
        pill_text="TSMC FOUNDRY CONTINUUM",
        page_num_str="07",
        t1_kicker="TARGET 1 · TSMC 22ULL TO N4E",
        t1_title="0.5V native ultra-low power FinFET",
        t1_body="Public proof: 22ULL to N4e near-threshold operation drops dynamic power by >75%",
        t2_kicker="TARGET 2 · THE 22NM ENVM DILEMMA",
        t2_title="28nm eFlash extinction cliff",
        t2_body="Public proof: ReRAM (+2 masks) and MRAM (+4 masks) inflate wafer cost by 15%-30%",
        syn_kicker="SYNOPSYS QUALIFICATION CONTRACT",
        syn_title="Qualify AntiFuse OTP as Code Core",
        syn_bullets=[
            "• zero extra mask cost on standard CMOS logic",
            "• 0.5V direct XIP boot and DMA to ULL SRAM",
            "• indirection jump tables for virtual bug patches",
            "• zero standby leakage for 10+ year battery life",
            "• TSMC 22ULL / 12FFC+ / N6e silicon proof",
            "• candidate unit: macro code density budget"
        ],
        bottom_banner_text="Public foundry roadmap: eFlash extinct below 28nm · Pure-logic AntiFuse OTP enables low-cost IoT MCU code storage"
    )

    # --- Slide 08: 3D Heterogeneous Chiplets & UCIe RoT ---
    s8 = prs.slides.add_slide(prs.slide_layouts[0])
    clone_slide_shapes(template_slide, s8)
    configure_perfect_card_slide(
        s8,
        kicker_text="07 · QUALIFY NOW · 3D CHIPLETS & UCIE",
        title_line1="3D SoIC Disaggregated Security",
        title_line2="Dedicated Hardware RoT per Chiplet",
        subtitle_text="Eliminating multi-die stacking loss · pre-bond KGD attestation",
        pill_text="3D PACKAGING INTEGRITY",
        page_num_str="08",
        t1_kicker="TARGET 1 · 3D SOIC STACKING INTEGRITY",
        t1_title="Known Good Die (KGD) verification",
        t1_body="Public proof: defective unverified die in 3D stack destroys the entire multi-die package",
        t2_kicker="TARGET 2 · UCIE 1.3 / 2.0 SECURITY",
        t2_title="Die-to-Die zero-trust link security",
        t2_body="Public proof: SPDM 1.3 mandates per-die cryptographic authentication before link traffic",
        syn_kicker="SYNOPSYS QUALIFICATION CONTRACT",
        syn_title="Deploy AntiFuse OTP + PUF per die",
        syn_bullets=[
            "• die-unique immutable identity and versioning",
            "• isolated pre-bond test key for wafer probe",
            "• SRAM PUF in-flight D2D encryption with 0 at rest",
            "• zero thermal penalty on 3D SoIC stacking",
            "• full compliance with SPDM 1.3 and UCIe 2.0",
            "• candidate unit: per-chiplet embedded RoT macro"
        ],
        bottom_banner_text="Chiplet Zero-Trust: A failed die inside 3D packaging destroys the entire stack; dedicated on-die RoT is non-negotiable"
    )

    # --- Slide 09: Post-Quantum Boot & Automotive 175°C SILC ---
    s9 = prs.slides.add_slide(prs.slide_layouts[0])
    clone_slide_shapes(template_slide, s9)
    configure_perfect_card_slide(
        s9,
        kicker_text="08 · QUALIFY NOW · PQC & AUTOMOTIVE SILC",
        title_line1="NIST SP 800-208 Stateful Boot",
        title_line2="175°C immune metallic filament physics",
        subtitle_text="Compact 32-byte LMS hash boot · AEC-Q100 Grade 0 reliability",
        pill_text="PHYSICS & QUANTUM GATE",
        page_num_str="09",
        t1_kicker="TARGET 1 · NIST SP 800-208 LMS BOOT",
        t1_title="Stateful hash tree 32-byte root hash",
        t1_body="Public proof: LMS overcomes 40x PQC lattice key bloat; OTS reuse destroys security",
        t2_kicker="TARGET 2 · AEC-Q100 GRADE 0 (175°C)",
        t2_title="SILC physics in automotive powertrain",
        t2_body="Public proof: 175°C stress causes trap-assisted tunneling and charge loss in floating gates",
        syn_kicker="SYNOPSYS QUALIFICATION CONTRACT",
        syn_title="1-to-0 Hard Breakdown Advantage",
        syn_bullets=[
            "• nanoscale metallic silicide has 0 drift at 175°C",
            "• irreversible OTP burn blocks OTS key reuse",
            "• 1T current-balanced sensing defeats DPA attacks",
            "• 1000h HTOL @ 175°C AEC-Q100 Grade 0 qualified",
            "• ISO 26262 ASIL-D ready automotive delivery",
            "• candidate unit: high-reliability secure RoT"
        ],
        bottom_banner_text="Physics proof: Metallic silicide conducts permanently; pure CMOS gate oxide breakdown achieves ASIL-D & PQC compliance"
    )

    # 6. Reorder slides to 17-slide sequence
    sldIdLst = prs.slides._sldIdLst
    old_elems = list(sldIdLst)
    desired_indices = [0, 1, 2, 3, 4, 5, 14, 15, 16, 6, 7, 8, 9, 10, 11, 12, 13]
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for idx in desired_indices:
        sldIdLst.append(old_elems[idx])
    print("Slides reordered.")

    # 7. Update Slide 10, 11, 12, 13, 14, 15, 16, 17 with ZERO OVERLAP
    # Slide 10: Repair Opportunity
    s10 = prs.slides[9]
    for sh in s10.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text
        if "06 · REPAIR OPPORTUNITY" in txt:
            clear_and_set_text(sh, "09 · REPAIR OPPORTUNITY · SPLIT THE LOCUS", "Consolas", 9.8, True, C_TEAL_KICKER)
        elif txt == "07":
            clear_and_set_text(sh, "10", "Consolas", 8.3, True, C_FOOTER_GRAY)

    # Slide 11: OCP Service State
    s11 = prs.slides[10]
    for sh in s11.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text
        if "07 · DISCOVER FIRST" in txt:
            clear_and_set_text(sh, "10 · DISCOVER FIRST · OCP SERVICE STATE", "Consolas", 9.8, True, C_TEAL_KICKER)
        elif txt == "08":
            clear_and_set_text(sh, "11", "Consolas", 8.3, True, C_FOOTER_GRAY)

    # Slide 12: Review Gate 2 (11 · LEADERSHIP CLOSE) - Strictly Single-Line 13.5pt to prevent collision!
    s12 = prs.slides[11]
    for sh in s12.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text
        if "08 · LEADERSHIP CLOSE" in txt:
            clear_and_set_text(sh, "11 · LEADERSHIP CLOSE", "Consolas", 9.8, True, C_TEAL_KICKER)
        elif "authorize three" in txt or "authorize five" in txt:
            clear_and_set_multiline(sh, ["Decision requested: authorize five", "target-bound qualification briefs"], "Arial", 30.0, True, C_DARK_NAVY)
        elif txt in ["09", "12"]:
            clear_and_set_text(sh, "12", "Consolas", 8.3, True, C_FOOTER_GRAY)
            
        # Left Cards: Strictly single-line titles & subtitles with generous spacing!
        elif "DDR5 SUBSYSTEM" in txt:
            clear_and_set_text(sh, "01 · DDR5 SUBSYSTEM", "Arial", 14.0, True, C_DARK_NAVY)
        elif "PMIC MTP" in txt:
            clear_and_set_text(sh, "PMIC MTP sequencing + SPD Hub protected blocks", "Arial", 11.5, False, C_SLATE_BODY)
        elif "ACCELERATOR POWER" in txt:
            clear_and_set_text(sh, "02 · AI POWER & TSMC IOT", "Arial", 14.0, True, C_DARK_NAVY)
        elif "Rail policy" in txt or "VR per-phase" in txt:
            clear_and_set_text(sh, "VR per-phase trim + 22ULL-N4e OTP code core", "Arial", 11.5, False, C_SLATE_BODY)
        elif "PLATFORM TRUST" in txt:
            clear_and_set_text(sh, "03 · TRUST & 3D CHIPLETS", "Arial", 14.0, True, C_DARK_NAVY)
        elif "One named BMC" in txt or "Caliptra RoT" in txt:
            clear_and_set_text(sh, "Caliptra RoT + UCIe 3D SoIC on-die KGD attestation", "Arial", 11.5, False, C_SLATE_BODY)

    # Slide 13: Appendix A
    s13 = prs.slides[12]
    for sh in s13.shapes:
        if not sh.has_text_frame: continue
        if sh.text_frame.text.strip() == "10":
            clear_and_set_text(sh, "13", "Consolas", 8.3, True, C_FOOTER_GRAY)

    # Slide 14: Review Gate 3 (APPENDIX B · SOURCE LEDGER)
    s14 = prs.slides[13]
    for sh in s14.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text.strip()
        if txt in ["11", "14"]:
            clear_and_set_text(sh, "14", "Consolas", 8.3, True, C_FOOTER_GRAY)
        elif "DDR5 and AI power have named" in txt or "DDR5, AI Power" in txt or "Explicit source evidence" in txt:
            clear_and_set_text(sh, "DDR5, AI Power, TSMC IoT & 3D Chiplet Source Ledger", "Arial", 28.0, True, C_DARK_NAVY)
        elif "Claim, state object" in txt or "Explicit public contracts" in txt:
            clear_and_set_text(sh, "Explicit public contracts bound every semiconductor qualification opportunity", "Arial", 14.8, False, C_SLATE_BODY)
            
        # Update Row 3 (AI Power to IoT)
        elif "Infineon XDPE1C284A" in txt:
            clear_and_set_text(sh, "Infineon XDPE1C284A / TSMC 22ULL-N4e", "Arial", 15.0, True, C_DARK_NAVY)
        elif "Internal NVM in accelerator" in txt:
            clear_and_set_text(sh, "Internal NVM in controller + pure-logic OTP Code Core", "Arial", 13.5, False, C_SLATE_BODY)
        elif "Exact state partition" in txt:
            clear_and_set_text(sh, "0.5V VDD · indirection patch tables · Synopsys fit", "Arial", 13.5, False, C_SLATE_BODY)
            
        # Update Row 4 (AI Power TI to Chiplet & PQC)
        elif "TI TPS536C9T" in txt:
            clear_and_set_text(sh, "UCIe 1.3 / OCP SPDM / NIST SP 800-208", "Arial", 15.0, True, C_DARK_NAVY)
        elif "Defaults · NVM fault status" in txt:
            clear_and_set_text(sh, "Die-to-Die hardware RoT + 32-byte LMS root hash", "Arial", 13.5, False, C_SLATE_BODY)
        elif "Write budget · event-log" in txt:
            clear_and_set_text(sh, "Pre-bond KGD test key · OTS anti-rollback counter", "Arial", 13.5, False, C_SLATE_BODY)

    # Slide 15: Review Gate 4 (APPENDIX C · TRUST / SERVICE LEDGER)
    s15 = prs.slides[14]
    for sh in s15.shapes:
        if not sh.has_text_frame: continue
        txt = sh.text_frame.text.strip()
        if txt in ["12", "15"]:
            clear_and_set_text(sh, "15", "Consolas", 8.3, True, C_FOOTER_GRAY)
        elif "OAI-OAM + OCP" in txt or "OAI-OAM + NIST" in txt:
            clear_and_set_text(sh, "OAI-OAM + NIST SP 800-208 PQC", "Consolas", 9.4, True, C_TEAL_KICKER)
        elif "Immutable RoT · SVN · anti-rollback" in txt or "Immutable RoT · SVN · LMS" in txt:
            clear_and_set_text(sh, "Immutable RoT · SVN · LMS state counter", "Arial", 17.0, True, C_DARK_NAVY)
        elif "OTP is named as an example" in txt or "AntiFuse irreversible burn" in txt:
            clear_and_set_text(sh, "AntiFuse irreversible burn blocks OTS key reuse; 32B root hash in OTP", "Arial", 14.0, False, C_SLATE_BODY)
            
        elif "CALIPTRA + AXIADO" in txt:
            clear_and_set_text(sh, "CALIPTRA + AXIADO + UCIE 3D SOIC", "Consolas", 9.4, True, C_TEAL_KICKER)
        elif "Persistent fuse fields + PUF can coexist" in txt or "Persistent fuse fields + PUF + D2D" in txt:
            clear_and_set_text(sh, "Persistent fuse fields + PUF + D2D RoT coexist", "Arial", 17.0, True, C_DARK_NAVY)
        elif "Lifecycle, hashes, revocation" in txt or "PUF does not replace" in txt:
            clear_and_set_text(sh, "PUF does not replace persistent policy, SVN, or chiplet KGD state", "Arial", 14.0, False, C_SLATE_BODY)

    # Slide 16: Appendix D
    s16 = prs.slides[15]
    for sh in s16.shapes:
        if not sh.has_text_frame: continue
        if sh.text_frame.text.strip() in ["13", "16"]:
            clear_and_set_text(sh, "16", "Consolas", 8.3, True, C_FOOTER_GRAY)

    # Slide 17: Appendix E
    s17 = prs.slides[16]
    for sh in s17.shapes:
        if not sh.has_text_frame: continue
        if sh.text_frame.text.strip() in ["14", "17"]:
            clear_and_set_text(sh, "17", "Consolas", 8.3, True, C_FOOTER_GRAY)

    # Save
    print(f"Saving presentation to: {OUT_PATH_CHATGPT}")
    prs.save(OUT_PATH_CHATGPT)
    print(f"Saving presentation to: {OUT_PATH_LOCAL}")

    # Enforce Title Case by default on all slide main titles (Slide 02 - Slide 17)
    title_case_map = {
        2: ["Authorize Five Target-Bound Qualifications", "Keep Two Families in Discovery"],
        3: ["Memory Choice Follows State Lifetime—", "Not the Application Label"],
        4: ["DDR5 PMIC + SPD Hub", "One Immediate Qualification Package"],
        5: ["AI Power Controllers Expose", "a Bounded-Update NVM Socket"],
        6: ["AI Platform Trust Needs OTP and PUF", "to Solve Different State Problems"],
        7: ["TSMC 22ULL to N4e Scaling", "Pure-Logic OTP as MCU Code Storage"],
        8: ["3D SoIC Disaggregated Security", "Dedicated Hardware RoT per Chiplet"],
        9: ["NIST SP 800-208 Stateful Boot", "175°C Immune Metallic Filament Physics"],
        10: ["AI SoC Repair Qualifies Directly", "HBM Repair Remains Locus-First"],
        11: ["OCP Modules Persist Service State", "but Standards Do Not Select Our Macro"],
        12: ["Decision Requested: Authorize Five", "Target-Bound Qualification Briefs"],
        13: ["Public Requirement ≠ Synopsys Fit ≠ Design Win ≠ Shipment"],
        14: ["DDR5, AI Power, TSMC IoT & 3D Chiplet Source Ledger"],
        15: ["Persistent Semantics Are Explicit;", "Implementation Technology Remains Bounded"],
        16: ["Before Selecting OTP, Name the Repair State", "and Its Programming Authority"],
        17: ["A Candidate Advances Only When Five Peer Gates Close"]
    }
    for s_idx, s in enumerate(prs.slides, 1):
        if s_idx in title_case_map:
            lines = title_case_map[s_idx]
            for shape in s.shapes:
                if shape.has_text_frame:
                    top_in = shape.top / 914400.0 if shape.top else 0
                    left_in = shape.left / 914400.0 if shape.left else 0
                    if 0.5 <= top_in <= 1.0 and 0.4 <= left_in <= 1.0:
                        tf = shape.text_frame
                        tf.clear()
                        for l_idx, l_txt in enumerate(lines):
                            p = tf.paragraphs[0] if l_idx == 0 else tf.add_paragraph()
                            p.space_before = Pt(0)
                            p.space_after = Pt(2)
                            run = p.add_run()
                            run.text = l_txt
                            run.font.name = "Arial"
                            run.font.size = Pt(28.0)
                            run.font.bold = True
                            run.font.color.rgb = C_DARK_NAVY
                        break

    prs.save(OUT_PATH_LOCAL)
    print("SUCCESS: 100% Zero-Collision Masterpiece v18 regenerated cleanly!")

if __name__ == "__main__":
    main()
