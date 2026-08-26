import os
import sys
import pptx

pptx_path = r"C:/Users/Sam/Documents/ChatGPT/NVM/Secure_Storage_SRAM_PUF_TSMC_OIP_Draft_v7_Visual_Enhanced.pptx"

if not os.path.exists(pptx_path):
    print(f"File not found: {pptx_path}")
    sys.exit(1)

prs = pptx.Presentation(pptx_path)
print(f"Presentation loaded. Total slides: {len(prs.slides)}")
print(f"Slide width: {prs.slide_width}, Slide height: {prs.slide_height}")

for idx, slide in enumerate(prs.slides):
    print(f"\n" + "="*50)
    print(f"SLIDE {idx+1}")
    print("="*50)
    for shape_idx, shape in enumerate(slide.shapes):
        shape_type = type(shape).__name__
        name = shape.name
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        
        if shape.has_text_frame:
            texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if texts:
                print(f"  [Shape {shape_idx}: {name} ({shape_type}) at ({left},{top})]:")
                for t in texts:
                    print(f"    - {t}")
        elif shape.has_table:
            print(f"  [Table {shape_idx}: {name}] ({len(shape.table.rows)}x{len(shape.table.columns)}):")
            for r_idx, row in enumerate(shape.table.rows):
                row_cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                print(f"    Row {r_idx}: {' | '.join(row_cells)}")
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            print(f"  [Picture {shape_idx}: {name} at ({left},{top}) dim=({width}x{height})]")
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            print(f"  [Group {shape_idx}: {name} at ({left},{top})]")
        else:
            print(f"  [Shape {shape_idx}: {name} ({shape_type})]")
