import os
import pptx

pptx_path = r"C:/Users/Sam/Documents/ChatGPT/NVM/Secure_Storage_SRAM_PUF_TSMC_OIP_Draft_v7_Visual_Enhanced.pptx"
prs = pptx.Presentation(pptx_path)

os.makedirs("extracted_images", exist_ok=True)

for i, slide in enumerate(prs.slides):
    for j, shape in enumerate(slide.shapes):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_bytes = image.blob
            ext = image.ext
            filename = f"extracted_images/slide_{i+1}_shape_{j}.{ext}"
            with open(filename, "wb") as f:
                f.write(image_bytes)
            print(f"Slide {i+1}, Shape {j}: Saved {filename} ({len(image_bytes)} bytes, format: {ext})")
