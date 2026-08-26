import pptx

prs = pptx.Presentation("Secure_Storage_SRAM_PUF_Mitigating_OTP_Leakage_v1.pptx")
print(f"Verified PPTX. Total slides: {len(prs.slides)}")
print(f"Slide width: {prs.slide_width.inches:.3f} in, Slide height: {prs.slide_height.inches:.3f} in")

for i, slide in enumerate(prs.slides):
    pictures = [s for s in slide.shapes if s.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE]
    text_boxes = [s for s in slide.shapes if s.has_text_frame]
    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else "No notes"
    print(f"Slide {i+1:2d}: {len(text_boxes)} text blocks, {len(pictures)} pictures, Notes length: {len(notes)} chars")
