import os
import pptx

pptx_path = r"C:/Users/Sam/Documents/ChatGPT/NVM/Secure_Storage_SRAM_PUF_TSMC_OIP_Draft_v7_Visual_Enhanced.pptx"
prs = pptx.Presentation(pptx_path)

output = []
for idx, slide in enumerate(prs.slides):
    output.append(f"==================================================")
    output.append(f"SLIDE {idx+1}")
    output.append(f"==================================================")
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if texts:
                output.append(f"  [Text {shape_idx}]: " + " // ".join(texts))
        elif shape.has_table:
            output.append(f"  [Table {shape_idx}]:")
            for r_idx, row in enumerate(shape.table.rows):
                row_cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                output.append(f"    Row {r_idx}: " + " | ".join(row_cells))
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            output.append(f"  [Picture {shape_idx}]")

with open("slides_dump.txt", "w", encoding="utf-8") as f:
    f.write("\n".join(output))

print("Dumped", len(prs.slides), "slides to slides_dump.txt")
