"""Module for generating the 100% English TSMC OIP Secure Storage Presentation Deck.

Embeds stunning, photorealistic 3D AI-generated semiconductor visuals on every slide,
with 100% English slide copy and comprehensive Bilingual (EN/ZH) Speaker Notes.
"""

import os
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Output configuration
OUTPUT_PPTX = "Secure_Storage_SRAM_PUF_Mitigating_OTP_Leakage_v1.pptx"
DIAGRAMS_DIR = "generated_diagrams"

# Color Palette Constants
C_DARK_BG = RGBColor(6, 27, 41)       # #061B29
C_CARD_BG = RGBColor(11, 37, 56)      # #0B2538
C_CARD_BORDER = RGBColor(30, 71, 101) # #1E4765
C_CYAN = RGBColor(0, 240, 255)        # #00F0FF
C_AMBER = RGBColor(255, 140, 66)      # #FF8C42
C_EMERALD = RGBColor(16, 185, 129)    # #10B981
C_PURPLE = RGBColor(168, 85, 247)     # #A855F7
C_RED = RGBColor(239, 68, 68)         # #EF4444
C_WHITE = RGBColor(255, 255, 255)     # #FFFFFF
C_MUTED = RGBColor(148, 163, 184)     # #94A3B8
C_LIGHT_BLUE = RGBColor(56, 189, 248) # #38BDF8


def apply_dark_background(slide) -> None:
    """Apply deep navy background to slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK_BG


def add_header(slide, category: str, title: str, subtitle: str, slide_num_str: str) -> None:
    """Add standard executive slide header in 100% English."""
    # Category Eyebrow
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.5), Inches(0.35))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_CYAN

    # Title
    tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(10.5), Inches(0.6))
    tf2 = tx_box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE

    # Subtitle
    if subtitle:
        tx_box3 = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(10.5), Inches(0.35))
        tf3 = tx_box3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = subtitle
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = C_MUTED

    # Slide Number Badge (Top Right)
    num_box = slide.shapes.add_textbox(Inches(11.8), Inches(0.45), Inches(0.9), Inches(0.4))
    tf_n = num_box.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.text = slide_num_str
    p_n.alignment = PP_ALIGN.RIGHT
    p_n.font.size = Pt(11.5)
    p_n.font.bold = True
    p_n.font.color.rgb = C_CYAN


def add_footer(slide, footer_note: str = "") -> None:
    """Add standard executive slide footer in 100% English."""
    f_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(7.5), Inches(0.3))
    tf = f_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TSMC OIP TECHNICAL BRIEFING  |  SECURE STORAGE ARCHITECTURES  |  CONFIDENTIAL"
    p.font.size = Pt(8)
    p.font.color.rgb = C_MUTED

    if footer_note:
        r_box = slide.shapes.add_textbox(Inches(7.5), Inches(6.95), Inches(5.0), Inches(0.3))
        tf_r = r_box.text_frame
        p_r = tf_r.paragraphs[0]
        p_r.text = footer_note
        p_r.alignment = PP_ALIGN.RIGHT
        p_r.font.size = Pt(8)
        p_r.font.bold = True
        p_r.font.color.rgb = C_LIGHT_BLUE


def add_structured_card(slide, left: float, top: float, width: float, height: float,
                        title: str, bullets: List[str], border_color: RGBColor = C_CYAN) -> None:
    """Add a dark executive card with structured bullet points in English."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.2)

    tx_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.12), Inches(width - 0.3), Inches(height - 0.24))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    if title:
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = border_color
        p_t.space_after = Pt(4)

    for b in bullets:
        p = tf.add_paragraph() if title or len(tf.paragraphs) > 0 else tf.paragraphs[0]
        p.text = b
        p.font.size = Pt(9)
        p.font.color.rgb = C_WHITE
        p.space_after = Pt(3)


def add_thematic_diagram(slide, left: float, top: float, width: float, height: float, image_path: str) -> None:
    """Embed the enhanced photorealistic AI visual diagram."""
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        print(f"[Warning] Image path not found: {image_path}")


# ==============================================================================
# SLIDE BUILDERS (100% ENGLISH SLIDES + 3D AI VISUALS + BILINGUAL NOTES)
# ==============================================================================

def build_slide_1(prs: Presentation) -> None:
    """Slide 1: Cover Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)

    # Title Box
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.8), Inches(1.8))
    tf = t_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "TSMC OIP TECHNICAL BRIEFING · AUGUST 2026"
    p0.font.size = Pt(10.5)
    p0.font.bold = True
    p0.font.color.rgb = C_CYAN
    p0.space_after = Pt(5)

    p1 = tf.add_paragraph()
    p1.text = "Mitigating OTP Data Leakage Risks in Advanced SoCs"
    p1.font.size = Pt(23)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Eliminating Raw Data Exposure in Leading-Edge Silicon Through Secure Storage Architectures"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = C_LIGHT_BLUE
    p2.space_after = Pt(4)

    p3 = tf.add_paragraph()
    p3.text = "SRAM PUF (Dynamic Root)  +  AES-256 (Line-Speed Cipher)  +  AntiFuse OTP (Scrambled Matrix)  +  Secure Controller"
    p3.font.size = Pt(10)
    p3.font.color.rgb = C_MUTED

    # Center Visual Diagram: AI Hero 3D Render
    add_thematic_diagram(slide, left=0.8, top=2.45, width=11.733, height=4.25,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_01.jpg"))

    add_footer(slide, "Core Thesis: Decouple physical storage from usable cryptographic secrets.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "Good morning, esteemed TSMC OIP partners and semiconductor architects. Today, we address a critical security challenge facing advanced node SoCs: 'Mitigating OTP Data Leakage Risks in Advanced SoCs Through Secure Storage Architectures.'\n"
        "Recent public hacking challenges, such as the Raspberry Pi RP2350 disclosure, have demonstrated that conventional AntiFuse OTP memory bits are physically observable under voltage fault injection and Focused Ion Beam (FIB) / Passive Voltage Contrast (PVC) probing. When root keys or firmware secrets are directly stored in OTP, physical readout leads to complete IP theft and device cloning.\n"
        "Our Secure Storage solution integrates SRAM PUF, AES-256 hardware crypto engines, AntiFuse OTP, and a Secure Controller into a single pre-integrated IP subsystem, ensuring that root keys are dynamically generated at power-up and never stored on-chip, turning OTP into a cryptographically protected repository.\n\n"
        "[中文講者備忘與論述重點]\n"
        "各位 TSMC OIP 夥伴與晶片架構師好。今天報告的主題聚焦於先進製程 SoC 中的 OTP 資料外洩防護。\n"
        "從 Raspberry Pi RP2350 公開挑戰賽可知，傳統 AntiFuse OTP 記憶體在 FIB/PVC 與電壓故障注入下，其實體 0/1 位元已不再具備物理不透明性。若直接於 OTP 存放根金鑰，將面臨災難性的複製風險。\n"
        "本 Secure Storage 架構將 SRAM PUF（動態根金鑰）、AES-256（硬體全加密）、AntiFuse OTP（亂序密文）與安全控制器整合為單一硬體子系統，達成『永久保存資料，根金鑰斷電零駐留』的根本安全突破。"
    )


def build_slide_2(prs: Presentation) -> None:
    """Slide 2: The Executive Thesis - Power-off State."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "01 / The Executive Thesis", "Security is Decided When the System is Powered Off",
               "Durable Storage Decoupled from Permanent Root Key Persistence", "02 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="Power-Off State (At-Rest Security)",
                        bullets=[
                            "• SRAM PUF unpowered: Root key is physically ABSENT",
                            "• AntiFuse OTP retains only AES-256 scrambled ciphertext",
                            "• Invasive physical readout recovers zero usable root secrets",
                            "• Eliminates the static at-rest attack surface of NVM/OTP-PUFs"
                        ], border_color=C_RED)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Power-Up State (Active Runtime Reconstruction)",
                        bullets=[
                            "• Transistor mismatch dynamically reconstructs 256-bit key",
                            "• Root key restricted entirely within hardware trust boundary",
                            "• On-the-fly decryption buffer instantly zeroized after task",
                            "• Achieves complete separation of retention vs key residence"
                        ], border_color=C_CYAN)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_02.jpg"))

    add_footer(slide, "Design Rule: Durable storage + Non-persistent root key + Hardware reconstruction")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "The highest level of hardware security is determined by what remains on the silicon when the system is powered down.\n"
        "Traditional OTP solves the retention problem, but fails the confidentiality challenge under invasive physical inspection. In our Secure Storage architecture, when power is removed, the SRAM PUF cells retain zero charge and zero state—the root key simply does not exist. Only scrambled AES-256 ciphertext remains in OTP. At power-up, the device reconstructs its unique root key on the fly within the secure hardware boundary, ensuring zero permanent key exposure.\n\n"
        "[中文講者備忘與論述重點]\n"
        "硬體安全的最高防線取決於斷電時晶片中還留下什麼。\n"
        "傳統 OTP 僅保證了資料留存 (Retention)，但無法保證實體機密性 (Confidentiality)。Secure Storage 的核心設計在於：斷電時 SRAM PUF 無任何電荷與金鑰殘留，OTP 內僅有 AES-256 密文；上電時才在內部安全邊界重建金鑰並即時解密，徹底杜絕晶片靜態提取風險。"
    )


def build_slide_3(prs: Presentation) -> None:
    """Slide 3: Threat Landscape - RP2350 Physical Attacks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "02 / Threat Landscape", "Public Attacks Expose the Limits of Stored-Bit Secrecy",
               "Lessons from the RP2350 Challenge: Physical Probing & Fault Injection Bypass Software Locks", "03 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="RP2350 Attack Demonstrations",
                        bullets=[
                            "• Voltage Fault Injection (FI): Bypassed boot timing & lock bits",
                            "• FIB / Passive Voltage Contrast: Imaged antifuse breakdown state",
                            "• Proved: Bit-cell opacity cannot serve as final security barrier",
                            "• Plain OTP Model: Extracting physical bits equals full compromise"
                        ], border_color=C_RED)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Architectural Design Consequence",
                        bullets=[
                            "• Assume physical observability: Do not rely solely on cell opacity",
                            "• Cryptographic transformation: Turn readout into ciphertext problem",
                            "• Composed defense: Ephemeral PUF root + AES-256 + Scrambling",
                            "• Architectural Rule: OTP Physical Readout != Secret Recovery"
                        ], border_color=C_AMBER)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_03.jpg"))

    add_footer(slide, "Real-world Signal: Memory security is a system property, not merely a bit-cell property.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "The RP2350 Hacking Challenge results in early 2025 served as a critical wake-up call for the semiconductor industry.\n"
        "Security researchers (such as IOActive) demonstrated that standard failure analysis equipment—specifically Focused Ion Beam (FIB) and Passive Voltage Contrast (PVC)—can directly inspect the physical programmed state of antifuse cells. Combined with voltage glitching during boot to bypass software permission wrappers, the raw bits of plain OTP were fully recovered.\n"
        "This proves that bit-cell opacity is not a sufficient defense. Our design response is clear: we assume the physical layer can eventually be imaged, but ensure that any extracted bitstream is cryptographically uncrackable ciphertext.\n\n"
        "[中文講者備忘與論述重點]\n"
        "RP2350 挑戰賽證實：硬體逆向工程團隊可利用標準失效分析設備（FIB 與被動電壓對比 PVC），直接讀出 AntiFuse 擊穿位元。\n"
        "若加上開機時序的電壓微擾 (Voltage Glitch)，傳統的暫存器鎖定與軟體保護將被輕易繞過。因此，系統架構不能假設『位元永遠無法被讀出』，而必須以密碼學手段確保『讀出結果僅為隨機密文』。"
    )


def build_slide_4(prs: Presentation) -> None:
    """Slide 4: Defense in Depth - Multi-Layer Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "03 / Defense in Depth", "Multi-Layer Pipeline: Restricting Physical Recovery to Ciphertext",
               "Four Independent Protective Barriers from Asset Ingestion to Silicon Storage", "04 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="Target Sensitive Assets",
                        bullets=[
                            "• Firmware anti-rollback monotonic counters",
                            "• AI model weight decryption keys & proprietary neural IP",
                            "• Device-unique private identity keys (ECC / RSA)",
                            "• Secure boot parameters & authorization policies"
                        ], border_color=C_LIGHT_BLUE)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Four-Tier Shield Mechanism",
                        bullets=[
                            "• Tier 1: SRAM PUF Dynamic Root (Zero at-rest footprint)",
                            "• Tier 2: AES-256 Hardware Encryption (Ciphertext barrier)",
                            "• Tier 3: Pseudo-random Address Scrambler (Obscures mapping)",
                            "• Tier 4: Split-Channel 1T Twin-Cell OTP (Power-balanced read)"
                        ], border_color=C_EMERALD)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_04.jpg"))

    add_footer(slide, "Defense Principle: Address scrambling raises mapping cost; AES-256 creates the mathematical barrier.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "Here we illustrate the four-tier defense-in-depth pipeline.\n"
        "When high-value assets—such as AI model keys, firmware counters, or device identities—enter the subsystem, they are encrypted via AES-256 using keys derived on-the-fly from the SRAM PUF root. The resulting ciphertext is then scrambled across physical rows and columns before programming into the 1T antifuse array.\n"
        "This creates an impenetrable dual barrier: address scrambling dramatically elevates the reverse engineering cost of physical reconstruction, while AES-256 provides mathematical confidentiality against raw bit recovery.\n\n"
        "[中文講者備忘與論述重點]\n"
        "本頁闡述四層多維度縱深防禦管線。\n"
        "敏感資產（AI 權重金鑰、防回滾計數器、私鑰）在寫入前，先經由 PUF 衍生金鑰進行 AES-256 加密，再經由位址混淆器 (Address Scrambler) 打亂邏輯對實體陣列的映射，最終寫入 1T 差動 AntiFuse 陣列。位址混淆提高尋址逆向成本，AES-256 建立數學級防線。"
    )


def build_slide_5(prs: Presentation) -> None:
    """Slide 5: Integrated Architecture Subsystem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "04 / Integrated Architecture", "Unified Subsystem: Four Blocks Under One Trust Boundary",
               "Pre-Integrated Hardware Security Subsystem Eliminating Cross-IP Vulnerabilities", "05 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="Four Core Hardware Blocks",
                        bullets=[
                            "• SRAM PUF: Dynamic key generator & silicon fingerprint core",
                            "• Crypto Engine: High-throughput line-speed AES-256 core",
                            "• AntiFuse OTP: Persistent encrypted storage matrix",
                            "• Secure Controller: APB slave, address scrambler & access FSM"
                        ], border_color=C_CYAN)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Subsystem Integration Advantages",
                        bullets=[
                            "• Eliminates cross-IP interface timing & reset glitch vulnerabilities",
                            "• Standard AMBA APB interface simplifies SoC host integration",
                            "• Hardware-enforced region access control & tamper response",
                            "• Single accountable IP provider for full security closure"
                        ], border_color=C_PURPLE)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_05.jpg"))

    add_footer(slide, "Pre-integrated architecture eliminates cross-IP assumptions and hand-off risks.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "A common pitfall in SoC security is assembling discrete IP blocks from different vendors—purchasing PUF from Vendor A, Controller from Vendor B, and OTP from Vendor C. Such fragmented architectures create unverified hand-offs, bus probe exposures, and reset synchronization glitches.\n"
        "Our Secure Storage IP delivers a pre-integrated, monolithic hardware subsystem. Key generation, cryptographic acceleration, address scrambling, and NVM access are entirely enclosed within a certified internal trust boundary, presenting only a clean, policy-protected APB interface to the host SoC.\n\n"
        "[中文講者備忘與論述重點]\n"
        "晶片安全最常見的破口往往是不同供應商 IP 之間的串接縫隙（如重置競爭、匯流排探針暴露）。\n"
        "Secure Storage 提供四合一單一 IP 交付：SRAM PUF、AES-256 引擎、AntiFuse OTP 與安全控制器封閉於同一硬體信任邊界內，向外部 SoC 僅暴露受保護的標準 APB 介面，消弭跨區塊假設風險。"
    )


def build_slide_6(prs: Presentation) -> None:
    """Slide 6: Root-Key Lifecycle & Zeroization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "05 / Key Lifecycle", "Root-Key Lifecycle: Strict Residency Control & Instant Zeroization",
               "Minimizing Key Exposure Windows Through Hardware-Enforced State Transitions", "06 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="Four-Phase Key Lifecycle",
                        bullets=[
                            "• Phase 1 (Power-Off): Root key ABSENT; OTP holds ciphertext",
                            "• Phase 2 (Power-Up): SRAM mismatch measured & key reconstructed",
                            "• Phase 3 (Authorized Use): Session keys derived for APB crypto",
                            "• Phase 4 (Zeroization): Immediate register flush upon task end"
                        ], border_color=C_AMBER)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Zero-Residency Security Assurance",
                        bullets=[
                            "• Core Metric: Minimize key residency time across execution",
                            "• Hardware tamper triggers force single-cycle latch clear",
                            "• Volatile register state dissipates within milliseconds of power cut",
                            "• Fully compliant with stringent anti-tamper security standards"
                        ], border_color=C_EMERALD)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_06.jpg"))

    add_footer(slide, "Security Property: Minimize key residency, not merely key visibility.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "Key management philosophy must focus on 'Minimizing Key Residency' rather than merely attempting to conceal stored keys.\n"
        "Our hardware state machine ensures that root key material exists only during the brief execution window required for cryptographic operations. As soon as reading, writing, or authentication is complete—or if a physical tamper event is detected—the controller executes hardware zeroization, flushing all key latches to zero in a single clock cycle.\n\n"
        "[中文講者備忘與論述重點]\n"
        "密鑰管理的核心準則是『Minimize Key Residency（最小化金鑰駐留時間）』。\n"
        "透過硬體狀態機控管，根金鑰僅在執行加解密時暫態存在。一旦操作完成或觸發入侵微擾偵測，硬體控制器會在一週期內執行 Zeroization（強制清零），徹底消除側信道洩漏的時間窗口。"
    )


def build_slide_7(prs: Presentation) -> None:
    """Slide 7: SRAM PUF Reconstruction & Helper Data Security."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "06 / PUF Physics & Security", "SRAM PUF Reconstruction: Information-Theoretically Secure Helper Data",
               "Public Activation Code Discloses Mathematically Zero Shannon Information Regarding Root Key", "07 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="PUF Key Reconstruction Pipeline",
                        bullets=[
                            "• Physical source: MOSFET threshold mismatch in standard SRAM",
                            "• Noise correction: BCH Error Correction Code (ECC) engine",
                            "• Privacy amplification: Compresses entropy to remove bias",
                            "• Output: 100% stable, unique 256-bit cryptographic root key"
                        ], border_color=C_CYAN)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Activation Code Security Guarantees",
                        bullets=[
                            "• Activation Code contains only syndrome data for ECC alignment",
                            "• Mathematical proof: Leaks exactly 0 bits of mutual information",
                            "• Publicly stored in Flash/OTP without compromising root secret",
                            "• Rigorously audited by leading academic and commercial certifiers"
                        ], border_color=C_EMERALD)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_07.jpg"))

    add_footer(slide, "Information-Theoretic Proof: Public helper data discloses mathematically zero key bits.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "A common technical question is: 'Does storing the Activation Code (helper data) publicly create a new attack vector?'\n"
        "The answer is an unequivocal no. The fuzzy extractor is constructed with information-theoretic security guarantees. The Activation Code consists purely of error-correction syndrome offsets. In information theory terms, the mutual Shannon information between the public helper data and the derived 256-bit secret key is mathematically zero. An adversary with infinite computing power cannot extract the key from the helper data alone.\n\n"
        "[中文講者備忘與論述重點]\n"
        "針對客戶常見的疑問：『公開存放的 Activation Code（輔助資料）是否會洩漏金鑰？』\n"
        "SRAM PUF 模糊提取器 (Fuzzy Extractor) 具備資訊理論安全性 (Information-Theoretic Security)。Activation Code 僅記錄糾錯症狀碼 (Syndrome)，其與最終 256-bit 金鑰之互資訊量 (Mutual Information) 嚴格為零，無任何密鑰位元洩漏風險。"
    )


def build_slide_8(prs: Presentation) -> None:
    """Slide 8: Commercial Proof & Reliability Metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "07 / Commercial Evidence", "Commercial Proof: 1.5B+ Devices and Extreme Automotive Reliability",
               "Field Maturity Informs—and Target Silicon Evaluation Closes the Security Claim", "08 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="Global Production Deployment Track Record",
                        bullets=[
                            "• 1.5B+ devices shipped with SRAM PUF hardware technology",
                            "• 15+ years commercial deployment in aerospace, defense & auto",
                            "• Node-agnostic architecture validated from 350nm down to 2nm",
                            "• Pure standard CMOS compatible: No extra masks or custom layers"
                        ], border_color=C_CYAN)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Extreme Environment & Security Standards",
                        bullets=[
                            "• Qualified across -40°C to 150°C PVT & accelerated aging models",
                            "• AEC-Q100 Grade 1 automotive certification compliance",
                            "• PSA Certified Level 3 & SESIP Level 3 hardware security",
                            "• NIST CAVP cryptographic algorithm validation compliance"
                        ], border_color=C_AMBER)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_08.jpg"))

    add_footer(slide, "Field maturity reduces uncertainty; target silicon evaluation closes the claim.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "Security claims must be backed by massive production maturity and stringent third-party certifications.\n"
        "SRAM PUF technology has been deployed in more than 1.5 billion commercial devices over 15+ years, spanning defense, aerospace, automotive, and IoT. Operating reliably across -40°C to 150°C with zero key error rate (0 ppm KER), it holds PSA Certified Level 3 and SESIP Level 3 accreditations. Because it relies on standard CMOS mismatch, it ports seamlessly across process nodes without requiring exotic foundry steps.\n\n"
        "[中文講者備忘與論述重點]\n"
        "SRAM PUF 技術具備超過 15 億顆晶片量產實績與 15 年以上商用經驗（廣泛用於航太、國防與車用晶片）。\n"
        "在 -40°C 至 150°C 寬溫域與 15 年老化模型下維持 0 ppm 金鑰重建錯誤率。並通過 AEC-Q100 Grade 1、PSA Level 3、SESIP Level 3 與 NIST CAVP 國際權威認證。"
    )


def build_slide_9(prs: Presentation) -> None:
    """Slide 9: Integration is the Product - Single Contract."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "08 / Integrated Delivery", "Integration is the Product: Single-Provider Accountability",
               "Customers Buy One Accountable Outcome: Root, Crypto, Storage, Control & Lifecycle", "09 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="Five-Pillar Cohesive Delivery",
                        bullets=[
                            "• 1. ROOT: SRAM PUF dynamic key generation",
                            "• 2. CRYPTO: High-speed line-rate AES-256 core",
                            "• 3. STORAGE: AntiFuse OTP permanent ciphertext array",
                            "• 4. CONTROL: Secure APB controller & address scrambling",
                            "• 5. LIFECYCLE: Zero-trust provisioning & lifetime support"
                        ], border_color=C_PURPLE)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Eliminating Multi-Vendor Fragmentation Risks",
                        bullets=[
                            "• Multi-vendor risks: Incompatible resets, bus leaks, finger-pointing",
                            "• Monolithic IP: Single contract, pre-verified silicon boundary",
                            "• Accelerates SoC time-to-market and simplifies cert audit trails",
                            "• Comprehensive wafer test, production provisioning & response"
                        ], border_color=C_CYAN)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_09.jpg"))

    add_footer(slide, "A secure primitive can still fail when interfaces, resets or ownership are fragmented.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "Security is an integrated property, not a catalog of separate primitives.\n"
        "When SoC teams buy separate security blocks from disparate suppliers, integration bugs and unverified reset states create severe vulnerability windows. With our Secure Storage solution, customers license a unified five-pillar system (Root, Crypto, Storage, Control, Lifecycle) under a single contract with guaranteed silicon closure and dedicated qualification collateral.\n\n"
        "[中文講者備忘與論述重點]\n"
        "安全是系統整合的結果，而非零散元件的拼湊。\n"
        "若分別採購 PUF、控制器與 OTP，整合時的時序與重置漏洞往往成為攻擊目標。本方案以單一 IP 交付五大支柱（Root, Crypto, Storage, Control, Lifecycle），為客戶提供單一責任窗口與完整的流片前驗證保證。"
    )


def build_slide_10(prs: Presentation) -> None:
    """Slide 10: TSMC OIP Readiness & Advanced Nodes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "09 / TSMC OIP Ecosystem", "TSMC OIP Readiness: Leading Process Nodes & Evidence Closure",
               "Silicon-Verified Collateral Across N7, N6, N5, N4P, N3P and Automotive Variants", "10 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="TSMC Advanced Process Node Support",
                        bullets=[
                            "• HPC & AI Nodes: N7, N6, N5, N4P, N3P silicon correlation",
                            "• Automotive Grades: N7A, N5A AEC-Q100 Grade 1 qualified",
                            "• Inherent PUF portability: No process-specific analog tuning",
                            "• Ready for TSMC GAA (N2) next-generation roadmap"
                        ], border_color=C_CYAN)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="OIP Silicon Evidence Closure Matrix",
                        bullets=[
                            "• Full TSMC silicon correlation reports across PVT corners",
                            "• Fault Injection (FI) & Side-Channel Analysis (SCA) test data",
                            "• FIB / Passive Voltage Contrast physical characterization",
                            "• Complete design deliverables: GDSII, LEF, LIB, Verilog, Testbenches"
                        ], border_color=C_AMBER)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_10.jpg"))

    add_footer(slide, "TSMC OIP readiness: Close evidence loop with qualified macro and target SoC boundary.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "Within the TSMC Open Innovation Platform (OIP), we deliver comprehensive silicon correlation across all flagship nodes, including N7, N6, N5, N4P, and N3P, as well as automotive N7A and N5A variants.\n"
        "We adhere to the 'Evidence Closure' principle: we provide customers with exhaustive PVT characterization tables, fault injection resistance reports, and physical FIB/PVC analysis data, ensuring first-time-right and first-time-secure silicon execution.\n\n"
        "[中文講者備忘與論述重點]\n"
        "在 TSMC OIP 架構下，本方案全面支援 N7、N6、N5、N4P、N3P 及車規 N7A/N5A 節點。\n"
        "我們堅持『Evidence Closure（證據閉環）』原則：提供完整的矽關聯報告、故障注入 (FI) 與 FIB/PVC 實體測試數據，並交付完整 GDSII/LEF/LIB/Verilog 檔案，確保一次流片即成功。"
    )


def build_slide_11(prs: Presentation) -> None:
    """Slide 11: Executive Conclusion & Value Proposition."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "10 / Executive Conclusion", "Conclusion: Decoupling Physical Storage from Usable Secrets",
               "Delivering First-Time-Secure Silicon Assurance for AI, Automotive, IoT and HPC", "11 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="Core Paradigm Breakthrough",
                        bullets=[
                            "• Core Breakthrough: OTP Readout != Secret Recovery",
                            "• Permanent Retention: Scrambled AES-256 ciphertext in AntiFuse",
                            "• Ephemeral Root: SRAM PUF dynamically reconstructed at boot",
                            "• Invasive Resilience: Eliminates FIB, PVC & voltage glitch risks"
                        ], border_color=C_CYAN)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Strategic Business & Time-to-Market Value",
                        bullets=[
                            "• Drastically mitigates multi-million dollar post-silicon re-spin risk",
                            "• Accelerates time-to-market for mission-critical AI, Auto & HPC",
                            "• Single accountable IP contract simplifies security compliance",
                            "• Engage today for target node macro specs and evaluation kits"
                        ], border_color=C_EMERALD)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_11.jpg"))

    add_footer(slide, "Next Step: Agree the evidence plan for target node, macro configuration and SoC boundary.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "In conclusion, conventional plain OTP architectures can no longer safeguard sensitive secrets against modern invasive physical attacks.\n"
        "By fundamentally decoupling physical non-volatile storage from usable cryptographic root keys, Secure Storage converts OTP from an observable attack surface into a cryptographically protected repository. This guarantees first-time-secure silicon for AI, automotive, and HPC applications while protecting corporate IP and time-to-market. Thank you, and we welcome your questions.\n\n"
        "[中文講者備忘與論述重點]\n"
        "總結而言，傳統 Plain OTP 架構已無法抵禦現代侵入式實體攻擊。\n"
        "Secure Storage 透過『解耦物理儲存與可用秘密』，將 OTP 轉化為受密碼學保護的密文儲存庫，實現『OTP 讀出 ≠ 秘密外洩』。為 AI、車用、IoT 與 HPC 晶片提供 First-Time-Secure 矽智財保證，大幅降低流片失敗風險。"
    )


def build_slide_12(prs: Presentation) -> None:
    """Slide 12: Technical Appendix - RP2350 Deep Dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_dark_background(slide)
    add_header(slide, "Appendix / Threat Deep-Dive", "Technical Deep-Dive: RP2350 Attack Paths & Countermeasures",
               "Comprehensive Failure Mode Analysis and Secure Storage Architectural Responses", "12 / 12")

    add_structured_card(slide, left=0.8, top=1.75, width=4.5, height=2.4,
                        title="RP2350 Physical Attack Vector Analysis",
                        bullets=[
                            "• Path 1 (Fault Injection): Glitched VDD at boot to bypass lock checks",
                            "• Path 2 (Invasive Probing): FIB/PVC imaged antifuse breakdown bits",
                            "• Path 3 (Sensor Evasion): Shaped glitch transients bypassed detectors",
                            "• Path 4 (Impact): Plain OTP exposed plaintext keys & boot secrets"
                        ], border_color=C_RED)

    add_structured_card(slide, left=0.8, top=4.3, width=4.5, height=2.45,
                        title="Secure Storage Defense Closure",
                        bullets=[
                            "• Response 1: All OTP data is AES-256 encrypted; bypassing locks yields cipher",
                            "• Response 2: Pseudo-random address scrambling destroys bit regularity",
                            "• Response 3: Mathematical confidentiality replaces reliance on reactive sensors",
                            "• Response 4: SRAM PUF dynamic key vanishes at power-off, zero at-rest secret"
                        ], border_color=C_EMERALD)

    add_thematic_diagram(slide, left=5.5, top=1.75, width=7.0, height=5.0,
                         image_path=os.path.join(DIAGRAMS_DIR, "ai_visual_slide_12.jpg"))

    add_footer(slide, "Use in Q&A if asked about technical trigger and failure-analysis attack countermeasures.")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "================================================================================\n"
        "[ENGLISH SPEAKER SCRIPT]\n"
        "This technical appendix is reserved for deep-dive Q&A discussions regarding the exact physical mechanisms demonstrated in the RP2350 challenge.\n"
        "We dissect the four critical attack paths—voltage fault injection, invasive FIB/PVC imaging, sensor evasion, and permission collapse—and demonstrate how Secure Storage creates a closed-loop cryptographic defense against each vector. Even if every detector is bypassed and every bit is physically imaged, the adversary cannot extract the ephemeral SRAM PUF root key.\n\n"
        "[中文講者備忘與論述重點]\n"
        "本技術附錄供 Q&A 環節深入探討 RP2350 實體攻擊途徑與防禦閉環。\n"
        "分析涵蓋電壓故障注入 (FI)、FIB/PVC 被動電壓對比、感測器規避與權限崩塌四大路徑，說明 Secure Storage 如何以 AES-256 加密、位址混淆與動態無常駐 PUF 金鑰達成完整的防禦閉環。"
    )


def main() -> None:
    """Build the complete 100% English 12-slide presentation deck with AI visuals & bilingual notes."""
    print("Building TSMC OIP Secure Storage Presentation Deck with 3D Photorealistic AI Visuals...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)

    prs.save(OUTPUT_PPTX)
    print(f"[Success] Presentation saved successfully: {OUTPUT_PPTX}")
    print(f"Total slides created: {len(prs.slides)}")


if __name__ == "__main__":
    main()
